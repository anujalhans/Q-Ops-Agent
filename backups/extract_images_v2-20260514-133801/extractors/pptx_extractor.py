from __future__ import annotations

import base64
import io
from typing import Any

from pptx import Presentation

from extract_images_v2.extractors.table_utils import table_rows_to_markdown, table_rows_to_plain_text
from extract_images_v2.extractors.visual_detection import (
    annotate_visual_candidates,
    convert_office_bytes_to_pdf_bytes,
    limit_visual_candidates,
    render_pdf_bytes,
    resolve_vision_config,
)


PICTURE_SHAPE_TYPE = 13


def _extract_slide_text(slide) -> str:
    lines: list[str] = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text and shape.text.strip():
            lines.append(shape.text.strip())
    return "\n".join(lines)


def _extract_slide_notes(slide, *, file_name: str, slide_number: int) -> list[dict[str, Any]]:
    notes: list[dict[str, Any]] = []
    try:
        notes_text = slide.notes_slide.notes_text_frame.text
    except Exception:
        notes_text = ""

    if notes_text and notes_text.strip():
        notes.append(
            {
                "annotationId": f"{file_name}_slide{slide_number}_notes",
                "fileName": file_name,
                "pageNumber": slide_number,
                "type": "speaker-notes",
                "title": "speaker-notes",
                "subject": "speaker-notes",
                "content": notes_text.strip(),
                "source": "pptx-speaker-notes",
            }
        )
    return notes


def _extract_slide_links(slide, *, file_name: str, slide_number: int) -> list[dict[str, Any]]:
    links: list[dict[str, Any]] = []
    link_index = 0
    for shape in slide.shapes:
        try:
            address = shape.click_action.hyperlink.address
        except Exception:
            address = None
        if address:
            links.append(
                {
                    "linkId": f"{file_name}_slide{slide_number}_link{link_index}",
                    "fileName": file_name,
                    "pageNumber": slide_number,
                    "uri": address,
                    "targetPage": None,
                    "kind": "external",
                    "source": "pptx-link",
                }
            )
            link_index += 1
    return links


def _extract_slide_tables(slide, *, file_name: str, slide_number: int) -> list[dict[str, Any]]:
    tables: list[dict[str, Any]] = []
    for table_index, shape in enumerate(slide.shapes):
        if not getattr(shape, "has_table", False):
            continue
        rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
        markdown = table_rows_to_markdown(rows)
        plain_text = table_rows_to_plain_text(rows)
        if not markdown and not plain_text:
            continue
        tables.append(
            {
                "tableId": f"{file_name}_slide{slide_number}_table{table_index}",
                "fileName": file_name,
                "pageNumber": slide_number,
                "tableIndex": table_index,
                "rowCount": len(rows),
                "columnCount": max((len(row) for row in rows), default=0),
                "markdown": markdown,
                "text": plain_text,
                "source": "pptx-table",
            }
        )
    return tables


def _extract_slide_images(slide, *, file_name: str, slide_number: int, doc_type: str) -> list[dict[str, Any]]:
    images: list[dict[str, Any]] = []
    image_index = 0
    for shape in slide.shapes:
        if shape.shape_type != PICTURE_SHAPE_TYPE:
            continue
        extension = shape.image.ext
        images.append(
            {
                "fileName": f"{file_name}_slide{slide_number}_image{image_index}.{extension}",
                "imageId": f"{file_name}_slide{slide_number}_img{image_index}",
                "base64": base64.b64encode(shape.image.blob).decode(),
                "docType": f"{doc_type}-IMAGE" if doc_type != "UNKNOWN" else "IMAGE",
                "imageSource": "embedded-image",
                "pageNumber": slide_number,
                "visualReason": ["embedded_image"],
                "mimeType": f"image/{extension}",
            }
        )
        image_index += 1
    return images


def _detect_pptx_visual_indicators(slide) -> list[str]:
    reasons: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_chart", False):
            if "chart" not in reasons:
                reasons.append("chart")
            continue
        if getattr(shape, "has_table", False):
            if "table_shapes" not in reasons:
                reasons.append("table_shapes")
            continue
        if shape.shape_type == PICTURE_SHAPE_TYPE:
            continue
        if hasattr(shape, "text") and shape.text and shape.text.strip():
            continue
        if "diagram_shapes" not in reasons:
            reasons.append("diagram_shapes")
    return reasons


def extract_pptx_document(*, content: bytes, context: dict[str, Any]) -> dict[str, Any]:
    file_name = context["fileName"]
    doc_type = context["docType"]
    vision_config = resolve_vision_config(context)

    presentation = Presentation(io.BytesIO(content))
    raw_text_parts: list[str] = []
    tables: list[dict[str, Any]] = []
    image_candidates: list[dict[str, Any]] = []
    annotations: list[dict[str, Any]] = []
    links: list[dict[str, Any]] = []
    rendered_pages: list[dict[str, Any]] = []
    warnings: list[str] = []
    all_visual_reasons: set[str] = set()

    for slide_index, slide in enumerate(presentation.slides):
        slide_number = slide_index + 1
        slide_text = _extract_slide_text(slide)
        if slide_text:
            raw_text_parts.append(slide_text)
        tables.extend(_extract_slide_tables(slide, file_name=file_name, slide_number=slide_number))
        image_candidates.extend(_extract_slide_images(slide, file_name=file_name, slide_number=slide_number, doc_type=doc_type))
        annotations.extend(_extract_slide_notes(slide, file_name=file_name, slide_number=slide_number))
        links.extend(_extract_slide_links(slide, file_name=file_name, slide_number=slide_number))
        all_visual_reasons.update(_detect_pptx_visual_indicators(slide))

    if all_visual_reasons:
        pdf_bytes, warning = convert_office_bytes_to_pdf_bytes(content, suffix=".pptx")
        if warning:
            warnings.append(warning)
        if pdf_bytes:
            rendered_pages = render_pdf_bytes(
                pdf_bytes,
                file_name=file_name,
                doc_type=doc_type,
                visual_reason=sorted(all_visual_reasons),
                prefix="pptx_rendered_slide",
                dpi=vision_config["renderDpi"],
            )

    images = annotate_visual_candidates(image_candidates)
    images, deferred_embedded_images = limit_visual_candidates(
        images,
        max_items=vision_config["maxEmbeddedImagesPerDocument"],
        overflow_label=f"{file_name} embedded images",
        warnings=warnings,
    )
    rendered_pages, deferred_rendered_pages = limit_visual_candidates(
        rendered_pages,
        max_items=vision_config["maxRenderedPagesPerDocument"],
        overflow_label=f"{file_name} rendered slides",
        warnings=warnings,
    )

    return {
        "fileType": "pptx",
        "pageCount": len(presentation.slides),
        "rawText": "\n\n".join(raw_text_parts),
        "images": images,
        "renderedPages": rendered_pages,
        "tables": tables,
        "annotations": annotations,
        "links": links,
        "warnings": warnings,
        "extractionStats": {
            "slidesProcessed": len(presentation.slides),
            "embeddedImagesDetected": len(image_candidates),
            "embeddedImagesExtracted": len(images),
            "embeddedImagesDeferred": deferred_embedded_images,
            "tablesExtracted": len(tables),
            "speakerNotesExtracted": len(annotations),
            "linksExtracted": len(links),
            "renderedSlideCandidatesDetected": len(rendered_pages) + deferred_rendered_pages,
            "renderedSlidesGenerated": len(rendered_pages),
            "renderedSlidesDeferred": deferred_rendered_pages,
            "visualCandidatesDetected": len(image_candidates) + len(rendered_pages) + deferred_embedded_images + deferred_rendered_pages,
            "visualCandidatesReturned": len(images) + len(rendered_pages),
            "visualCandidatesDeferred": deferred_embedded_images + deferred_rendered_pages,
            "renderDpi": vision_config["renderDpi"],
        },
    }
