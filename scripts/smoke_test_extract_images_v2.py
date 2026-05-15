from __future__ import annotations

import asyncio
import csv
import json
from pathlib import Path
import sys
import zipfile

import fitz
from docx import Document
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches
from starlette.datastructures import UploadFile


ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from extract_images_v2.main import process_document_v2

OUTPUT_DIR = ROOT / "docs" / "test_data" / "extract_images_v2_smoke"
RESULTS_DIR = OUTPUT_DIR / "results"
SAMPLES_DIR = OUTPUT_DIR / "samples"
OMNICART_DIR = ROOT / "docs" / "test_data" / "omnicart_payment_gateway_project"
SOURCE_IMAGE = OMNICART_DIR / "ui_ux_jpg" / "01_home_search_results.jpg"
SOURCE_PDF = OMNICART_DIR / "pdf_documents" / "BRD_OmniCart_Payment_Gateway_Project.pdf"


def ensure_dirs() -> None:
    RESULTS_DIR.mkdir(parents=True, exist_ok=True)
    SAMPLES_DIR.mkdir(parents=True, exist_ok=True)


def add_hyperlink(paragraph, text: str, url: str) -> None:
    part = paragraph.part
    relation_id = part.relate_to(url, RT.HYPERLINK, is_external=True)
    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(qn("r:id"), relation_id)

    run = OxmlElement("w:r")
    run_properties = OxmlElement("w:rPr")
    run_style = OxmlElement("w:rStyle")
    run_style.set(qn("w:val"), "Hyperlink")
    run_properties.append(run_style)
    run.append(run_properties)

    text_element = OxmlElement("w:t")
    text_element.text = text
    run.append(text_element)
    hyperlink.append(run)
    paragraph._p.append(hyperlink)


def create_txt_sample(path: Path) -> None:
    path.write_text(
        "Grooming transcript extract\n"
        "Participants: Product Owner, QA Lead, Architect\n"
        "Decision: support partial capture, refunds, and 3DS fallback validation.\n",
        encoding="utf-8",
    )


def create_csv_sample(path: Path) -> None:
    with path.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.writer(handle)
        writer.writerow(["requirement_id", "module", "priority"])
        writer.writerow(["REQ-101", "checkout", "high"])
        writer.writerow(["REQ-205", "refunds", "medium"])


def create_docx_sample(path: Path) -> None:
    document = Document()
    document.add_heading("OmniCart Payment Gateway Functional Notes", level=1)
    document.add_paragraph("The checkout flow must support cards, UPI, wallets, and net banking.")
    document.add_paragraph("This paragraph references the payment orchestration sequence diagram.")
    paragraph = document.add_paragraph("Reference architecture: ")
    add_hyperlink(paragraph, "Payment API Contract", "https://example.com/payment-api")

    table = document.add_table(rows=3, cols=3)
    headers = ["Field", "Validation", "Source"]
    values = [
        ["card_number", "Luhn + BIN", "FRD"],
        ["payment_status", "authorized/captured/refunded", "HLD"],
    ]
    for index, header in enumerate(headers):
        table.cell(0, index).text = header
    for row_index, row in enumerate(values, start=1):
        for col_index, value in enumerate(row):
            table.cell(row_index, col_index).text = value

    document.add_picture(str(SOURCE_IMAGE), width=Inches(4.5))
    document.save(path)


def create_pptx_sample(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Payment Gateway Flow"

    text_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(3.8), Inches(1.2))
    text_box.text_frame.text = "Checkout validates fraud score before hitting the payment gateway."

    table_shape = slide.shapes.add_table(3, 3, Inches(0.6), Inches(2.1), Inches(4.4), Inches(1.6))
    table = table_shape.table
    table.cell(0, 0).text = "Step"
    table.cell(0, 1).text = "System"
    table.cell(0, 2).text = "Outcome"
    table.cell(1, 0).text = "1"
    table.cell(1, 1).text = "Risk Engine"
    table.cell(1, 2).text = "Allow"
    table.cell(2, 0).text = "2"
    table.cell(2, 1).text = "Gateway"
    table.cell(2, 2).text = "Authorize"

    slide.shapes.add_picture(str(SOURCE_IMAGE), Inches(5.4), Inches(1.1), width=Inches(3.5))

    arrow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.CHEVRON,
        Inches(4.2),
        Inches(4.1),
        Inches(1.3),
        Inches(0.6),
    )
    arrow.text = "Route"
    arrow.click_action.hyperlink.address = "https://example.com/payment-routing"

    chart_data = ChartData()
    chart_data.categories = ["Cards", "UPI", "Wallets"]
    chart_data.add_series("Volume", (60, 25, 15))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.6),
        Inches(4.3),
        Inches(3.4),
        Inches(2.2),
        chart_data,
    )

    presentation.save(path)


def create_pdf_sample(path: Path) -> None:
    document = fitz.open()
    page = document.new_page(width=842, height=595)
    page.insert_text((48, 48), "OmniCart Payment Gateway Architecture", fontsize=18)
    page.insert_text(
        (48, 82),
        "The gateway router supports retry, idempotency, partial capture, and reconciliation.",
        fontsize=11,
    )

    link_rect = fitz.Rect(48, 96, 270, 112)
    page.insert_text((48, 108), "Payment API Specification", fontsize=10)
    page.insert_link({"kind": fitz.LINK_URI, "from": link_rect, "uri": "https://example.com/payment-api"})

    page.add_text_annot((360, 70), "Sticky note: validate 3DS fallback and timeout handling.")

    table_origin_x = 48
    table_origin_y = 150
    column_widths = [160, 180, 160]
    row_height = 28
    rows = [
        ["Step", "Validation", "Owner"],
        ["Auth", "3DS + fraud check", "Gateway"],
        ["Capture", "Partial capture allowed", "OMS"],
        ["Refund", "Async webhook update", "Finance"],
    ]

    x = table_origin_x
    for width in column_widths:
        page.draw_line((x, table_origin_y), (x, table_origin_y + row_height * len(rows)), color=(0, 0, 0), width=1)
        x += width
    page.draw_line((x, table_origin_y), (x, table_origin_y + row_height * len(rows)), color=(0, 0, 0), width=1)

    for row_index in range(len(rows) + 1):
        y = table_origin_y + row_height * row_index
        page.draw_line((table_origin_x, y), (table_origin_x + sum(column_widths), y), color=(0, 0, 0), width=1)

    for row_index, row in enumerate(rows):
        for col_index, value in enumerate(row):
            cell_x = table_origin_x + sum(column_widths[:col_index]) + 6
            cell_y = table_origin_y + row_height * row_index + 18
            page.insert_text((cell_x, cell_y), value, fontsize=10)

    page.draw_rect(fitz.Rect(520, 150, 720, 210), color=(0.1, 0.3, 0.7), width=1.5)
    page.draw_line((620, 210), (620, 280), color=(0.1, 0.3, 0.7), width=2)
    page.draw_line((620, 280), (560, 330), color=(0.1, 0.3, 0.7), width=2)
    page.draw_line((620, 280), (680, 330), color=(0.1, 0.3, 0.7), width=2)
    page.insert_text((548, 190), "Gateway Router", fontsize=11)
    page.insert_text((525, 344), "Fraud Engine", fontsize=10)
    page.insert_text((652, 344), "Bank Adapter", fontsize=10)

    page.insert_image(fitz.Rect(500, 360, 780, 540), filename=str(SOURCE_IMAGE))
    document.save(path)
    document.close()


def create_docx_with_comments_warnings(path: Path) -> None:
    with zipfile.ZipFile(path, "a") as zip_file:
        comments_xml = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:comments xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:comment w:id="0" w:author="QA Lead">
    <w:p><w:r><w:t>Review negative paths for partial capture and expired authorization.</w:t></w:r></w:p>
  </w:comment>
</w:comments>
"""
        zip_file.writestr("word/comments.xml", comments_xml)


def generate_samples() -> dict[str, Path]:
    ensure_dirs()

    txt_path = SAMPLES_DIR / "grooming_transcript_sample.txt"
    csv_path = SAMPLES_DIR / "requirements_sample.csv"
    docx_path = SAMPLES_DIR / "omnicart_sample.docx"
    pptx_path = SAMPLES_DIR / "omnicart_sample.pptx"
    pdf_path = SAMPLES_DIR / "omnicart_visual_sample.pdf"

    create_txt_sample(txt_path)
    create_csv_sample(csv_path)
    create_docx_sample(docx_path)
    create_docx_with_comments_warnings(docx_path)
    create_pptx_sample(pptx_path)
    create_pdf_sample(pdf_path)

    return {
        "pdf_custom": pdf_path,
        "pdf_existing": SOURCE_PDF,
        "docx": docx_path,
        "pptx": pptx_path,
        "txt": txt_path,
        "csv": csv_path,
        "image": SOURCE_IMAGE,
    }


async def post_file(file_path: Path) -> dict:
    with file_path.open("rb") as handle:
        upload = UploadFile(filename=file_path.name, file=handle)
        response = await process_document_v2(
            file=upload,
            projectName="Extraction Smoke Project",
            status="processing",
            jobId="SMOKE-EXTRACT-V2",
            projectId="project-smoke",
            requestedBy="tester-smoke",
            settingsVersion="2",
        )
    if response.status_code >= 400:
        raise RuntimeError(f"{file_path.name} failed with status {response.status_code}: {response.body.decode('utf-8', errors='ignore')}")
    return json.loads(response.body.decode("utf-8"))


def summarize_payload(name: str, payload: dict) -> dict:
    return {
        "sample": name,
        "fileType": payload.get("fileType"),
        "docType": payload.get("docType"),
        "pageCount": payload.get("pageCount"),
        "contentMode": payload.get("contentMode"),
        "containsText": payload.get("containsText"),
        "containsImages": payload.get("containsImages"),
        "rawTextChars": len(payload.get("rawText", "")),
        "imageCount": payload.get("imageCount", 0),
        "tableCount": len(payload.get("tables", [])),
        "renderedPagesCount": len(payload.get("renderedPages", [])),
        "annotationCount": len(payload.get("annotations", [])),
        "linkCount": len(payload.get("links", [])),
        "warnings": payload.get("warnings", []),
        "extractionStats": payload.get("extractionStats", {}),
    }


def main() -> None:
    sample_files = generate_samples()
    summary_rows = []

    for sample_name, sample_path in sample_files.items():
        payload = asyncio.run(post_file(sample_path))
        output_path = RESULTS_DIR / f"{sample_name}.json"
        output_path.write_text(json.dumps(payload, indent=2), encoding="utf-8")
        summary_rows.append(summarize_payload(sample_name, payload))

    summary_path = RESULTS_DIR / "summary.json"
    summary_path.write_text(json.dumps(summary_rows, indent=2), encoding="utf-8")

    print(json.dumps(summary_rows, indent=2))


if __name__ == "__main__":
    main()
