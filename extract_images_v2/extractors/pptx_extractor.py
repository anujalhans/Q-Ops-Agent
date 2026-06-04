from __future__ import annotations

import base64
import io
from typing import Any

from pptx import Presentation

from extract_images_v2.extractors.table_utils import table_rows_to_markdown, table_rows_to_plain_text
from extract_images_v2.extractors.visual_detection import annotate_visual_candidate


PICTURE_SHAPE_TYPE = 13


def _extract_slide_tables(slide, *, file_name: str, slide_number: int, warnings: list[str]) -> list[dict[str, Any]]:
    tables: list[dict[str, Any]] = []
    try:
        for table_index, shape in enumerate(slide.shapes):
            if not getattr(shape, "has_table", False):
                continue
            rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
            markdown = table_rows_to_markdown(rows)
            plain_text = table_rows_to_plain_text(rows)
            if not markdown and not plain_text:
                continue
            tables.append(
                {
                    "tableId": f"{file_name}_slide{slide_number}_table{table_index}",
                    "fileName": file_name,
                    "pageNumber": slide_number,
                    "tableIndex": table_index,
                    "rowCount": len(rows),
                    "columnCount": max((len(row) for row in rows), default=0),
                    "markdown": markdown,
                    "text": plain_text,
                    "source": "pptx-table",
                }
            )
    except Exception as exc:
        warnings.append(f"PPTX table extraction skipped on slide {slide_number}: {exc}")
    return tables


def _extract_slide_notes(slide, *, file_name: str, slide_number: int, warnings: list[str]) -> list[dict[str, Any]]:
    notes: list[dict[str, Any]] = []
    try:
        notes_text = slide.notes_slide.notes_text_frame.text
    except Exception:
        notes_text = ""

    try:
        if notes_text and notes_text.strip():
            notes.append(
                {
                    "annotationId": f"{file_name}_slide{slide_number}_notes",
                    "fileName": file_name,
                    "pageNumber": slide_number,
                    "type": "speaker-notes",
                    "title": "speaker-notes",
                    "subject": "speaker-notes",
                    "content": notes_text.strip(),
                    "source": "pptx-speaker-notes",
                }
            )
    except Exception as exc:
        warnings.append(f"PPTX speaker-note extraction skipped on slide {slide_number}: {exc}")
    return notes


def _extract_slide_links(slide, *, file_name: str, slide_number: int, warnings: list[str]) -> list[dict[str, Any]]:
    links: list[dict[str, Any]] = []
    link_index = 0
    try:
        for shape in slide.shapes:
            try:
                address = shape.click_action.hyperlink.address
            except Exception:
                address = None
            if not address:
                continue
            links.append(
                {
                    "linkId": f"{file_name}_slide{slide_number}_link{link_index}",
                    "fileName": file_name,
                    "pageNumber": slide_number,
                    "uri": address,
                    "targetPage": None,
                    "kind": "external",
                    "source": "pptx-link",
                }
            )
            link_index += 1
    except Exception as exc:
        warnings.append(f"PPTX link extraction skipped on slide {slide_number}: {exc}")
    return links


def _build_pptx_render_candidate(
    *,
    slide,
    file_name: str,
    doc_type: str,
    slide_index: int,
    table_count: int,
    annotation_count: int,
    link_count: int,
    image_count: int,
) -> dict[str, Any] | None:
    slide_number = slide_index + 1
    reasons: list[str] = []

    if table_count:
        reasons.append("table_shapes")
    if annotation_count:
        reasons.append("annotations")
    if link_count:
        reasons.append("document_links")
    if image_count:
        reasons.append("embedded_image")

    chart_count = 0
    vector_shape_count = 0
    text_length = 0
    for shape in slide.shapes:
        if getattr(shape, "has_chart", False):
            chart_count += 1
        if getattr(shape, "shape_type", None) != PICTURE_SHAPE_TYPE and not getattr(shape, "has_table", False):
            vector_shape_count += 1
        if hasattr(shape, "text") and shape.text:
            text_length += len(shape.text.strip())

    if chart_count:
        reasons.append("chart")
    if vector_shape_count >= 6:
        reasons.append("diagram_shapes")
    if text_length < 120 and (image_count or chart_count or table_count or vector_shape_count >= 6):
        reasons.append("low_text_density")

    meaningful_reasons = [reason for reason in reasons if reason != "document_links"]
    if not meaningful_reasons:
        return None

    return annotate_visual_candidate(
        {
            "fileName": f"{file_name}_render_candidate_slide{slide_index}.png",
            "imageId": f"{file_name}_render_candidate_slide{slide_index}",
            "docType": f"{doc_type}-IMAGE" if doc_type != "UNKNOWN" else "IMAGE",
            "imageSource": "rendered-page-candidate",
            "pageNumber": slide_number,
            "visualReason": reasons,
            "mimeType": "image/png",
        }
    )


def extract_pptx_document(*, content: bytes, context: dict[str, Any]) -> dict[str, Any]:
    file_name = context["fileName"]
    doc_type = context["docType"]
    extraction_config = context.get("extractionConfig", {})
    extract_tables = bool(extraction_config.get("extractTables", True))
    extract_annotations = bool(extraction_config.get("extractAnnotations", True))
    extract_links = bool(extraction_config.get("extractLinks", True))
    detect_rendered_pages = bool(extraction_config.get("detectRenderedPages", False))
    render_pages = bool(extraction_config.get("renderPages", False))

    presentation = Presentation(io.BytesIO(content))
    raw_text_parts: list[str] = []
    images: list[dict[str, Any]] = []
    tables: list[dict[str, Any]] = []
    annotations: list[dict[str, Any]] = []
    links: list[dict[str, Any]] = []
    visual_candidates: list[dict[str, Any]] = []
    warnings: list[str] = []

    for slide_index, slide in enumerate(presentation.slides):
        slide_number = slide_index + 1
        slide_tables: list[dict[str, Any]] = []
        slide_annotations: list[dict[str, Any]] = []
        slide_links: list[dict[str, Any]] = []
        if extract_tables:
            slide_tables = _extract_slide_tables(slide, file_name=file_name, slide_number=slide_number, warnings=warnings)
            tables.extend(slide_tables)
        if extract_annotations:
            slide_annotations = _extract_slide_notes(slide, file_name=file_name, slide_number=slide_number, warnings=warnings)
            annotations.extend(slide_annotations)
        if extract_links:
            slide_links = _extract_slide_links(slide, file_name=file_name, slide_number=slide_number, warnings=warnings)
            links.extend(slide_links)

        slide_image_count = 0
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                raw_text_parts.append(shape.text.strip())

            if shape.shape_type != PICTURE_SHAPE_TYPE:
                continue

            slide_image_count += 1
            extension = shape.image.ext
            images.append(
                {
                    "fileName": f"{file_name}_slide{slide_index}.{extension}",
                    "imageId": f"{file_name}_slide{slide_index}_img",
                    "base64": base64.b64encode(shape.image.blob).decode(),
                    "docType": f"{doc_type}-IMAGE" if doc_type != "UNKNOWN" else "IMAGE",
                    "imageSource": "embedded-image",
                    "pageNumber": slide_number,
                    "visualReason": ["embedded_image"],
                    "mimeType": f"image/{extension}",
                }
            )

        if detect_rendered_pages:
            candidate = _build_pptx_render_candidate(
                slide=slide,
                file_name=file_name,
                doc_type=doc_type,
                slide_index=slide_index,
                table_count=len(slide_tables),
                annotation_count=len(slide_annotations),
                link_count=len(slide_links),
                image_count=slide_image_count,
            )
            if candidate:
                visual_candidates.append(candidate)

    return {
        "fileType": "pptx",
        "pageCount": len(presentation.slides),
        "rawText": "\n".join(raw_text_parts),
        "images": images,
        "renderedPages": [],
        "visualCandidates": visual_candidates,
        "tables": tables,
        "annotations": annotations,
        "links": links,
        "warnings": warnings,
        "extractionStats": {
            "compatibilityMode": False,
            "safeTextRichMode": True,
            "renderingEnabled": False,
            "renderPagesRequested": render_pages,
            "renderDetectionEnabled": detect_rendered_pages,
            "extractTablesEnabled": extract_tables,
            "extractAnnotationsEnabled": extract_annotations,
            "extractLinksEnabled": extract_links,
            "slidesProcessed": len(presentation.slides),
            "embeddedImagesExtracted": len(images),
            "renderedPagesGenerated": 0,
            "visualCandidatesDetected": len(visual_candidates),
            "visualCandidatePages": ", ".join(str(item.get("pageNumber")) for item in visual_candidates),
            "tablesExtracted": len(tables),
            "speakerNotesExtracted": len(annotations),
            "annotationsExtracted": len(annotations),
            "linksExtracted": len(links),
        },
    }
